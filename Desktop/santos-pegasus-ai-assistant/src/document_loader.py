from __future__ import annotations

import csv
import json
import re
from datetime import datetime
from pathlib import Path
from typing import Iterable

import pandas as pd
from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from .models import DocumentChunk

SUPPORTED_EXTENSIONS = {
    ".pdf", ".csv", ".json", ".md", ".markdown", ".html", ".htm",
    ".txt", ".docx", ".xlsx", ".pptx"
}


def _category_from_path(path: Path) -> str:
    stem = path.stem.lower()
    mapping = {
        "onboarding": "Recursos Humanos",
        "backend": "Ingeniería Backend",
        "frontend": "Ingeniería Frontend",
        "incidente": "Operaciones y Calidad",
        "arquitectura": "Arquitectura y Sistemas",
        "directorio": "Comunicación Interna",
        "seguridad": "Legal y Compliance",
        "precios": "Financiero y Comercial",
    }
    for key, value in mapping.items():
        if key in stem:
            return value
    return "General"


def _mtime(path: Path) -> str:
    return datetime.fromtimestamp(path.stat().st_mtime).isoformat(timespec="seconds")


def load_pdf(path: Path) -> list[dict]:
    reader = PdfReader(str(path))
    sections = []
    for page_number, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            sections.append({"text": text, "location": f"Página {page_number}"})
    return sections


def load_csv_file(path: Path) -> list[dict]:
    df = pd.read_csv(path)
    sections = []
    for index, row in df.fillna("").iterrows():
        content = "; ".join(f"{column}: {row[column]}" for column in df.columns)
        sections.append({"text": content, "location": f"Fila {index + 2}"})
    return sections


def load_json_file(path: Path) -> list[dict]:
    data = json.loads(path.read_text(encoding="utf-8"))
    records = data if isinstance(data, list) else [data]
    return [
        {
            "text": json.dumps(record, ensure_ascii=False, indent=2),
            "location": f"Registro {i}",
        }
        for i, record in enumerate(records, start=1)
    ]


def load_markdown_file(path: Path) -> list[dict]:
    raw = path.read_text(encoding="utf-8")
    try:
        from markdown import markdown
        html = markdown(raw)
        text = BeautifulSoup(html, "html.parser").get_text("\n")
    except ImportError:
        # Fallback suficiente para conservar el contenido antes de instalar extras.
        text = re.sub(r"[`*_>#-]+", " ", raw)
    return [{"text": text, "location": "Documento completo"}]


def load_html_file(path: Path) -> list[dict]:
    soup = BeautifulSoup(path.read_text(encoding="utf-8"), "html.parser")
    return [{"text": soup.get_text("\n"), "location": "Documento completo"}]


def load_text_file(path: Path) -> list[dict]:
    return [{"text": path.read_text(encoding="utf-8"), "location": "Documento completo"}]


def load_docx_file(path: Path) -> list[dict]:
    doc = Document(path)
    sections, current_heading, buffer = [], "Inicio", []

    def flush():
        nonlocal buffer
        text = "\n".join(buffer).strip()
        if text:
            sections.append({"text": text, "location": f"Sección: {current_heading}"})
        buffer = []

    for paragraph in doc.paragraphs:
        value = paragraph.text.strip()
        if not value:
            continue
        if paragraph.style and paragraph.style.name.startswith("Heading"):
            flush()
            current_heading = value
        else:
            buffer.append(value)
    flush()
    return sections


def load_xlsx_file(path: Path) -> list[dict]:
    workbook = load_workbook(path, read_only=True, data_only=True)
    sections = []
    for sheet in workbook.worksheets:
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue
        headers = [str(v) if v is not None else f"columna_{i+1}" for i, v in enumerate(rows[0])]
        for row_number, row in enumerate(rows[1:], start=2):
            text = "; ".join(
                f"{headers[i]}: {value}" for i, value in enumerate(row) if value is not None
            )
            if text:
                sections.append({
                    "text": text,
                    "location": f"Hoja {sheet.title}, fila {row_number}"
                })
    return sections


def load_pptx_file(path: Path) -> list[dict]:
    presentation = Presentation(path)
    sections = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            sections.append({
                "text": "\n".join(texts),
                "location": f"Diapositiva {slide_number}"
            })
    return sections


LOADERS = {
    ".pdf": load_pdf,
    ".csv": load_csv_file,
    ".json": load_json_file,
    ".md": load_markdown_file,
    ".markdown": load_markdown_file,
    ".html": load_html_file,
    ".htm": load_html_file,
    ".txt": load_text_file,
    ".docx": load_docx_file,
    ".xlsx": load_xlsx_file,
    ".pptx": load_pptx_file,
}


def discover_documents(directory: Path) -> list[Path]:
    if not directory.exists():
        return []
    return sorted(
        p for p in directory.rglob("*")
        if p.is_file() and p.suffix.lower() in SUPPORTED_EXTENSIONS
    )


def load_document_sections(path: Path) -> list[dict]:
    loader = LOADERS.get(path.suffix.lower())
    if not loader:
        raise ValueError(f"Formato no soportado: {path.suffix}")
    return loader(path)


def load_all_sections(directory: Path) -> list[dict]:
    output = []
    for path in discover_documents(directory):
        category = _category_from_path(path)
        try:
            sections = load_document_sections(path)
        except Exception as exc:
            print(f"[WARN] No se pudo procesar {path.name}: {exc}")
            continue
        for section in sections:
            output.append({
                **section,
                "source": path.name,
                "category": category,
                "modified_at": _mtime(path),
                "owner": "Santos Pegasus Soluciones",
            })
    return output
